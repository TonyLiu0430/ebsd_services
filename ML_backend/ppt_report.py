from __future__ import annotations

import io
import math
import statistics
from datetime import datetime
from typing import Any, Dict, Iterable, List, Optional, Tuple

from matplotlib.backends.backend_agg import FigureCanvasAgg
from matplotlib.figure import Figure
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


COL_KEYS = ["C", "M", "E"]
ROW_KEYS = ["U", "M", "B"]
COL_LABELS = {"C": "內 Center", "M": "中 Middle", "E": "外 Edge"}
ROW_LABELS = {"U": "上 Up", "M": "中 Middle", "B": "下 Bottom"}
ROW_SERIES_INFO = [
    {"rowKey": "U", "color": RGBColor(239, 68, 68), "label": "上 (U)"},
    {"rowKey": "M", "color": RGBColor(59, 130, 246), "label": "中 (M)"},
    {"rowKey": "B", "color": RGBColor(16, 185, 129), "label": "下 (B)"},
]
ORIENT_LINE_SERIES_INFO = [
    {"orientationKey": "001", "index": 0, "color": RGBColor(64, 158, 255), "label": "[001]"},
    {"orientationKey": "011", "index": 1, "color": RGBColor(103, 194, 58), "label": "[011]"},
    {"orientationKey": "111", "index": 2, "color": RGBColor(230, 162, 60), "label": "[111]"},
]
ORIENT_LABELS = ["001", "011", "111"]
ORIENT_LINE_X_TICKS = [
    {"colKey": "C", "label": "內"},
    {"colKey": "M", "label": "中"},
    {"colKey": "E", "label": "外"},
]
FONT_FACE = "Microsoft JhengHei"
SUCCESS_BG = RGBColor(240, 249, 235)
SUCCESS_BORDER = RGBColor(179, 225, 157)
SUCCESS_TEXT = RGBColor(103, 194, 58)
PRIMARY_BG = RGBColor(236, 245, 255)
PRIMARY_BORDER = RGBColor(179, 216, 255)
PRIMARY_TEXT = RGBColor(64, 158, 255)
WARNING_BG = RGBColor(253, 246, 236)
WARNING_BORDER = RGBColor(243, 209, 158)
WARNING_TEXT = RGBColor(230, 162, 60)


FeatureMap = Dict[str, Any]
Snapshot = Dict[str, FeatureMap]
ReportData = Dict[str, Snapshot]


def build_pptx_report(
    *,
    selected_sample: str,
    golden_sample: str,
    selected_version_label: str,
    golden_version_label: str,
    report_data: ReportData,
    ipf_images: Dict[str, bytes],
    ipf_legend: Optional[bytes] = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    selected_snapshot = report_data.get(selected_sample, {})
    golden_snapshot = report_data.get(golden_sample, {})
    selected_positions = sorted(selected_snapshot.keys(), key=_sort_position_key)
    golden_positions = sorted(golden_snapshot.keys(), key=_sort_position_key)
    missing_golden = [pos for pos in selected_positions if pos not in golden_snapshot]
    missing_selected = [pos for pos in golden_positions if pos not in selected_snapshot]

    _add_cover_slide(prs, selected_sample, golden_sample, selected_version_label, golden_version_label)
    _add_grain_distribution_slide(
        prs,
        selected_snapshot,
        golden_snapshot,
        selected_version_label,
        golden_version_label,
    )
    if ipf_images:
        _add_ipf_slide(prs, selected_sample, selected_version_label, ipf_images, ipf_legend)
    _add_nine_grid_slide(prs, selected_sample, selected_version_label, selected_snapshot, missing_golden, missing_selected)
    _add_orientation_triangle_slide(prs, selected_sample, golden_sample, selected_version_label, golden_version_label, report_data, "20%")
    _add_orientation_triangle_slide(prs, selected_sample, golden_sample, selected_version_label, golden_version_label, report_data, "15%")
    _add_orientation_line_slide(prs, selected_sample, golden_sample, selected_version_label, golden_version_label, selected_snapshot, golden_snapshot, "20%")
    _add_orientation_line_slide(prs, selected_sample, golden_sample, selected_version_label, golden_version_label, selected_snapshot, golden_snapshot, "15%")

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = RGBColor(17, 24, 39),
    align: Optional[int] = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FACE
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _set_cell_text(cell, text: str, *, size: int = 9, bold: bool = False, color: RGBColor = RGBColor(31, 41, 55)):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FACE
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _style_header_cell(cell, text: str):
    cell.fill.solid()
    cell.fill.fore_color.rgb = SUCCESS_BG
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    _set_cell_text(cell, text, size=9, bold=True, color=SUCCESS_TEXT)


def _style_body_cell(cell, text: str, *, size: int = 10):
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    _set_cell_text(cell, text, size=size, color=RGBColor(31, 41, 55))


def _add_section_title(slide, text: str):
    _add_text(slide, text, 0.45, 0.28, 12.2, 0.38, size=20, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.76), Inches(12.35), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(229, 231, 235)
    line.line.fill.background()


def _add_picture_contain(slide, image: io.BytesIO, x: float, y: float, w: float, h: float):
    image.seek(0)
    with Image.open(image) as img:
        img_w, img_h = img.size
    image.seek(0)
    ratio = min(w / max(img_w, 1), h / max(img_h, 1))
    draw_w = img_w * ratio
    draw_h = img_h * ratio
    slide.shapes.add_picture(
        image,
        Inches(x + (w - draw_w) / 2),
        Inches(y + (h - draw_h) / 2),
        Inches(draw_w),
        Inches(draw_h),
    )


def _add_cover_slide(prs: Presentation, selected_sample: str, golden_sample: str, selected_version_label: str, golden_version_label: str):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    _add_text(slide, "EBSD 自動化報表", 0.75, 1.05, 11.8, 0.7, size=34, bold=True)
    _add_text(slide, f"分析樣本：{selected_sample} · {selected_version_label}", 0.8, 2.15, 8.5, 0.35, size=18)
    _add_text(slide, f"Golden 樣本：{golden_sample} · {golden_version_label}", 0.8, 2.65, 8.5, 0.35, size=18)
    _add_text(slide, "報告內容：Grain 粒徑分布、IPF 晶粒取向分佈圖、九宮格完整數據、晶粒取向比例與折線圖。", 0.8, 3.35, 11.7, 0.7, size=15)
    _add_text(slide, f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}", 0.82, 6.65, 5.8, 0.25, size=10, color=RGBColor(107, 114, 128))


def _add_grain_distribution_slide(
    prs: Presentation,
    selected_snapshot: Snapshot,
    golden_snapshot: Snapshot,
    selected_version_label: str,
    golden_version_label: str,
):
    values: List[float] = []
    for snapshot in (selected_snapshot, golden_snapshot):
        for features in snapshot.values():
            values.extend(_finite_list(features.get("grains", [])))
    positive = [v for v in values if v > 0]
    hist_bin_min = math.floor(min(positive) * 10) / 10 if positive else 7.1
    x_display_min = 0.0
    x_max = 200.0

    modes = [
        ("cdf", "Grain 粒徑分布比對 - 累積曲線"),
        ("hist", "Grain 粒徑分布比對 - 分布直方圖"),
        ("areaHist", "Grain 粒徑分布比對 - 面積加權直方圖"),
    ]

    chart_w = 3.72
    chart_h = 1.68
    x0 = 0.55
    y0 = 1.35
    for mode, title in modes:
        slide = _blank_slide(prs)
        _add_section_title(slide, title)
        _add_legend(slide, 0.65, 0.93, [(selected_version_label, RGBColor(59, 130, 246)), (f"{golden_version_label} (黃金)", RGBColor(245, 158, 11))])
        if mode != "cdf":
            _add_text(slide, f"Bin 起點 {hist_bin_min:.1f} um / Bin 數 16 / 顯示軸從 0 到 200 um", 6.2, 0.93, 6.2, 0.22, size=9, color=RGBColor(75, 85, 99))

        for r, row_key in enumerate(ROW_KEYS):
            for c, col_key in enumerate(COL_KEYS):
                pos = f"{col_key}-{row_key}"
                x = x0 + c * 4.15
                y = y0 + r * 1.86
                _draw_grain_chart(
                    slide,
                    x,
                    y,
                    chart_w,
                    chart_h,
                    pos,
                    _finite_list(selected_snapshot.get(pos, {}).get("grains", [])),
                    _finite_list(golden_snapshot.get(pos, {}).get("grains", [])),
                    mode,
                    x_display_min,
                    x_max,
                    hist_bin_min,
                )


def _add_ipf_slide(prs: Presentation, selected_sample: str, selected_version_label: str, ipf_images: Dict[str, bytes], ipf_legend: Optional[bytes]):
    slide = _blank_slide(prs)
    _add_section_title(slide, f"IPF 晶粒取向分佈圖 - {selected_version_label or selected_sample}")
    if ipf_legend:
        _add_ipf_legend(slide, ipf_legend, 11.25, 0.18, 0.85, 0.62)
    cell_w = 3.35
    cell_h = 1.78
    x0 = 1.25
    y0 = 1.13
    for r, row_key in enumerate(ROW_KEYS):
        for c, col_key in enumerate(COL_KEYS):
            pos = f"{col_key}-{row_key}"
            x = x0 + c * 3.85
            y = y0 + r * 1.95
            _add_text(slide, pos, x, y - 0.14, 0.8, 0.16, size=8, bold=True)
            data = ipf_images.get(pos)
            if data:
                try:
                    _add_picture_contain(slide, io.BytesIO(data), x, y + 0.05, cell_w, cell_h)
                except Exception:
                    _add_text(slide, "圖片載入失敗", x, y + 0.62, cell_w, 0.25, size=10, align=PP_ALIGN.CENTER)
            else:
                _add_placeholder(slide, x, y + 0.05, cell_w, cell_h, "未上傳")


def _add_ipf_legend(slide, ipf_legend: bytes, x: float, y: float, w: float, h: float):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.07), Inches(y - 0.05), Inches(w + 0.14), Inches(h + 0.24))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.color.rgb = RGBColor(220, 223, 230)
    _add_text(slide, "IPF 色鍵", x, y - 0.01, w, 0.14, size=6, bold=True, align=PP_ALIGN.CENTER)
    _add_picture_contain(slide, io.BytesIO(ipf_legend), x + 0.05, y + 0.11, w - 0.1, h - 0.08)
    _add_text(slide, "[111]", x + w * 0.36, y + 0.08, w * 0.28, 0.1, size=4, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "[001]", x - 0.02, y + h - 0.03, w * 0.36, 0.1, size=4, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "[101]", x + w * 0.68, y + h - 0.03, w * 0.36, 0.1, size=4, bold=True, align=PP_ALIGN.CENTER)


def _add_nine_grid_slide(
    prs: Presentation,
    selected_sample: str,
    selected_version_label: str,
    selected_snapshot: Snapshot,
    missing_golden: Iterable[str],
    missing_selected: Iterable[str],
):
    missing_golden_set = set(missing_golden)
    missing_selected_set = set(missing_selected)

    slide = _blank_slide(prs)
    _add_section_title(slide, f"九宮格完整數據 - {selected_version_label or selected_sample}")
    x0 = 1.1
    y0 = 1.38
    cell_w = 3.58
    cell_h = 1.58
    gap_x = 0.28
    gap_y = 0.18

    for c, col_key in enumerate(COL_KEYS):
        x = x0 + c * (cell_w + gap_x)
        _add_text(slide, COL_LABELS[col_key], x, 0.96, cell_w, 0.28, size=11, bold=True, color=SUCCESS_TEXT, align=PP_ALIGN.CENTER)
    for r, row_key in enumerate(ROW_KEYS):
        y = y0 + r * (cell_h + gap_y)
        _add_text(slide, ROW_LABELS[row_key], 0.35, y + 0.58, 0.62, 0.22, size=10, bold=True, color=RGBColor(75, 85, 99), align=PP_ALIGN.CENTER)
        for c, col_key in enumerate(COL_KEYS):
            pos = f"{col_key}-{row_key}"
            x = x0 + c * (cell_w + gap_x)
            _draw_nine_grid_data_card(slide, x, y, cell_w, cell_h, pos, selected_snapshot.get(pos), _position_status(pos, selected_snapshot.get(pos), missing_golden_set, missing_selected_set))


def _draw_nine_grid_data_card(slide, x: float, y: float, w: float, h: float, pos: str, features: Optional[FeatureMap], status: str):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = SUCCESS_BORDER if status == "OK" else WARNING_BORDER
    card.line.width = Pt(0.9)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.23))
    band.fill.solid()
    band.fill.fore_color.rgb = SUCCESS_BG if status == "OK" else WARNING_BG
    band.line.fill.background()
    _add_text(slide, pos, x + 0.08, y + 0.03, 0.55, 0.14, size=8, bold=True)
    _add_text(slide, status, x + w - 0.8, y + 0.03, 0.68, 0.14, size=7, bold=True, color=SUCCESS_TEXT if status == "OK" else WARNING_TEXT, align=PP_ALIGN.RIGHT)

    mean, max_v, p75, count = _grain_stats(features)
    orient20 = _orient_values(features, "20%") or [0.0, 0.0, 0.0]
    orient15 = _orient_values(features, "15%") or [0.0, 0.0, 0.0]

    _add_text(slide, "Grain Size (um)", x + 0.12, y + 0.32, 1.25, 0.14, size=7, bold=True, color=RGBColor(75, 85, 99))
    _add_text(slide, f"Mean  {mean}\nMax   {max_v}\nP75   {p75}\nCount {count}", x + 0.12, y + 0.5, 1.22, 0.64, size=7)

    _add_text(slide, "Orientation Ratio", x + 1.5, y + 0.32, 1.55, 0.14, size=7, bold=True, color=RGBColor(75, 85, 99))
    orient_table = slide.shapes.add_table(3, 4, Inches(x + 1.48), Inches(y + 0.51), Inches(w - 1.62), Inches(0.62)).table
    for idx, width in enumerate([0.34, 0.54, 0.54, 0.54]):
        orient_table.columns[idx].width = Inches(width)
    for col, label in enumerate(["", "[001]", "[011]", "[111]"]):
        _style_header_cell(orient_table.cell(0, col), label)
        _set_cell_text(orient_table.cell(0, col), label, size=5, bold=True, color=SUCCESS_TEXT)
    for row, (dev, values) in enumerate([("20deg", orient20), ("15deg", orient15)], start=1):
        _set_cell_text(orient_table.cell(row, 0), dev, size=5, bold=True, color=RGBColor(75, 85, 99))
        for col in range(3):
            _set_cell_text(orient_table.cell(row, col + 1), f"{values[col]:.1f}%", size=5)


def _draw_orientation_ratio_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    pos: str,
    selected_features: Optional[FeatureMap],
    golden_features: Optional[FeatureMap],
    dev: str,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = SUCCESS_BORDER
    card.line.width = Pt(0.9)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.24))
    band.fill.solid()
    band.fill.fore_color.rgb = SUCCESS_BG
    band.line.fill.background()
    _add_text(slide, pos, x + 0.1, y + 0.04, 0.55, 0.14, size=8, bold=True)
    _add_text(slide, f"Misorientation {dev}", x + w - 1.45, y + 0.04, 1.32, 0.14, size=7, bold=True, color=SUCCESS_TEXT, align=PP_ALIGN.RIGHT)

    selected_values = _orient_values(selected_features, dev) or [0.0, 0.0, 0.0]
    golden_values = _orient_values(golden_features, dev) or [0.0, 0.0, 0.0]
    table = slide.shapes.add_table(3, 4, Inches(x + 0.18), Inches(y + 0.42), Inches(w - 0.36), Inches(0.78)).table
    for idx, width in enumerate([0.72, 0.75, 0.75, 0.75]):
        table.columns[idx].width = Inches(width)
    for col, label in enumerate(["", "[001]", "[011]", "[111]"]):
        _style_header_cell(table.cell(0, col), label)
        _set_cell_text(table.cell(0, col), label, size=6, bold=True, color=SUCCESS_TEXT)
    row_specs = [
        ("Sample", selected_values, PRIMARY_BG, PRIMARY_TEXT),
        ("Golden", golden_values, WARNING_BG, WARNING_TEXT),
    ]
    for row, (label, values, bg, color) in enumerate(row_specs, start=1):
        table.cell(row, 0).fill.solid()
        table.cell(row, 0).fill.fore_color.rgb = bg
        _set_cell_text(table.cell(row, 0), label, size=6, bold=True, color=color)
        for col in range(3):
            _style_body_cell(table.cell(row, col + 1), f"{values[col]:.1f}%", size=7)


def _add_orientation_triangle_slide(
    prs: Presentation,
    selected_sample: str,
    golden_sample: str,
    selected_version_label: str,
    golden_version_label: str,
    report_data: ReportData,
    dev: str,
):
    slide = _blank_slide(prs)
    _add_section_title(slide, f"晶粒取向比例 - Misorientation {dev}")
    _add_text(slide, f"{selected_sample} · {selected_version_label} / {golden_sample} · {golden_version_label}", 0.65, 0.9, 11.8, 0.25, size=11, color=RGBColor(75, 85, 99))
    selected_snapshot = report_data.get(selected_sample, {})
    golden_snapshot = report_data.get(golden_sample, {})
    _add_triangle_slide_legend(slide, 0.7, 1.15, selected_version_label, golden_version_label)
    x0 = 0.6
    y0 = 1.55
    panel_w = 3.95
    panel_h = 5.25
    for c, col_key in enumerate(COL_KEYS):
        _draw_triangle_panel(slide, x0 + c * 4.25, y0, panel_w, panel_h, col_key, dev, selected_snapshot, golden_snapshot)


def _add_orientation_line_slide(
    prs: Presentation,
    selected_sample: str,
    golden_sample: str,
    selected_version_label: str,
    golden_version_label: str,
    selected_snapshot: Snapshot,
    golden_snapshot: Snapshot,
    dev: str,
):
    slide = _blank_slide(prs)
    _add_section_title(slide, f"晶粒取向折線圖 - Misorientation {dev}")
    _add_text(slide, f"{selected_sample} · {selected_version_label} / {golden_sample} · {golden_version_label}", 0.65, 0.9, 6.5, 0.25, size=11, color=RGBColor(75, 85, 99))
    _add_line_chart_legend(slide, 7.25, 0.9, selected_version_label, golden_version_label)
    chart_w = 3.42
    chart_h = 1.48
    x0 = 0.65
    y0 = 1.22
    for r, row_key in enumerate(ROW_KEYS):
        for c, col_key in enumerate(COL_KEYS):
            pos = f"{col_key}-{row_key}"
            x = x0 + c * 4.18
            y = y0 + r * 1.82
            _draw_position_orientation_line_chart(slide, x, y, chart_w, chart_h, pos, selected_snapshot.get(pos), golden_snapshot.get(pos), dev)


def _draw_grain_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    sample: List[float],
    golden: List[float],
    mode: str,
    x_display_min: float,
    x_max: float,
    hist_bin_min: float,
):
    _add_text(slide, title, x, y, w, 0.18, size=9, bold=True, align=PP_ALIGN.CENTER)
    chart_png = _render_grain_chart_png(sample, golden, mode, x_display_min, x_max, hist_bin_min)
    _add_picture_contain(slide, chart_png, x + 0.05, y + 0.22, w - 0.1, h - 0.25)


def _draw_position_orientation_line_chart(slide, x: float, y: float, w: float, h: float, pos: str, features: Optional[FeatureMap], golden_features: Optional[FeatureMap], dev: str):
    _add_text(slide, pos, x, y, w, 0.18, size=9, bold=True, align=PP_ALIGN.CENTER)
    chart_png = _render_position_orientation_chart_png(
        _orient_values(features, dev) or [0.0, 0.0, 0.0],
        _orient_values(golden_features, dev) or [0.0, 0.0, 0.0],
    )
    _add_picture_contain(slide, chart_png, x + 0.03, y + 0.2, w - 0.06, h - 0.22)


def _add_line_chart_legend(slide, x: float, y: float, selected_version_label: str, golden_version_label: str):
    sample_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y + 0.12), Inches(x + 0.3), Inches(y + 0.12))
    sample_line.line.color.rgb = PRIMARY_TEXT
    sample_line.line.width = Pt(2)
    _add_text(slide, f"分析 {selected_version_label}", x + 0.36, y, 1.6, 0.2, size=8, color=RGBColor(75, 85, 99))
    golden_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.05), Inches(y + 0.12), Inches(x + 2.35), Inches(y + 0.12))
    golden_line.line.color.rgb = WARNING_TEXT
    golden_line.line.width = Pt(2)
    golden_line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _add_text(slide, f"Golden {golden_version_label}", x + 2.41, y, 1.65, 0.2, size=8, color=RGBColor(75, 85, 99))


def _render_grain_chart_png(sample: List[float], golden: List[float], mode: str, x_display_min: float, x_max: float, hist_bin_min: float) -> io.BytesIO:
    fig = Figure(figsize=(3.15, 2.08), dpi=220)
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    if mode == "cdf":
        def plot(values: List[float], color: str, linestyle: str):
            xs, ys = _cdf_points(values)
            if xs:
                ax.plot(xs, ys, color=color, linewidth=1.8, linestyle=linestyle, solid_joinstyle="round")

        plot(golden, "#f59e0b", "--")
        plot(sample, "#3b82f6", "-")
        ax.set_ylim(0, 100)
        ax.set_yticks([0, 25, 50, 75, 100])
        ax.set_ylabel("Cumulative %", fontsize=7, color="#6b7280", labelpad=1)
        ax.set_xticks([0, 50, 100, 150, 200])
    else:
        area_weighted = mode == "areaHist"
        bins = 16
        bin_width = max(x_max - hist_bin_min, 1e-6) / bins
        edges = [hist_bin_min + bin_width * i for i in range(bins + 1)]
        sample_series = _histogram_series(sample, bins, hist_bin_min, x_max, area_weighted)
        golden_series = _histogram_series(golden, bins, hist_bin_min, x_max, area_weighted)
        ymax = max(1.0, *(sample_series + golden_series))
        y_step = max(1, math.ceil(ymax / 4))
        y_axis_max = y_step * 4
        ax.bar(
            edges[:-1],
            golden_series,
            width=bin_width,
            align="edge",
            facecolor=(245 / 255, 158 / 255, 11 / 255, 0.20),
            edgecolor=(245 / 255, 158 / 255, 11 / 255, 0.95),
            linewidth=1.0,
            linestyle=(0, (2.2, 1.4)),
            zorder=2,
        )
        ax.bar(
            edges[:-1],
            sample_series,
            width=bin_width,
            align="edge",
            facecolor=(15 / 255, 118 / 255, 110 / 255, 0.24),
            edgecolor=(15 / 255, 118 / 255, 110 / 255, 0.95),
            linewidth=1.0,
            zorder=3,
        )
        ax.set_ylim(0, y_axis_max)
        ax.set_yticks([0, y_axis_max * 0.25, y_axis_max * 0.5, y_axis_max * 0.75, y_axis_max])
        ax.set_ylabel("Area %" if area_weighted else "Counts", fontsize=7, color="#6b7280", labelpad=1)
        ax.set_xticks([0, *edges])
        ax.set_xticklabels([_format_axis_tick(v) for v in [0, *edges]], rotation=55, ha="right", fontsize=5.2, color="#6b7280")

    ax.set_xlim(x_display_min, x_max)
    ax.yaxis.grid(True, color="#f0f0f0", linewidth=0.7, zorder=0)
    ax.xaxis.grid(False)
    ax.tick_params(labelsize=6, colors="#6b7280", length=2)
    ax.set_xlabel("Grain (um)", fontsize=7, color="#6b7280", labelpad=1)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#9ca3af")
    ax.spines["bottom"].set_color("#9ca3af")
    ax.spines["left"].set_linewidth(0.75)
    ax.spines["bottom"].set_linewidth(0.75)
    fig.tight_layout(pad=0.25)

    output = io.BytesIO()
    canvas.print_png(output)
    output.seek(0)
    return output


def _render_position_orientation_chart_png(values: List[float], golden_values: List[float]) -> io.BytesIO:
    fig = Figure(figsize=(3.15, 1.25), dpi=220)
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    xs = [0, 1, 2]
    safe_values = [min(100.0, max(0.0, float(v))) for v in values]
    safe_golden_values = [min(100.0, max(0.0, float(v))) for v in golden_values]
    ax.plot(xs, safe_golden_values, color="#e6a23c", linewidth=1.8, marker="o", markersize=3.2, linestyle="--")
    ax.plot(xs, safe_values, color="#409eff", linewidth=1.8, marker="o", markersize=3.5)
    ax.set_xlim(-0.15, 2.15)
    ax.set_ylim(0, 100)
    ax.set_xticks(xs)
    ax.set_xticklabels(["[001]", "[011]", "[111]"], fontsize=6, color="#6b7280")
    ax.set_yticks([0, 25, 50, 75, 100])
    ax.tick_params(labelsize=6, colors="#6b7280", length=2)
    ax.yaxis.grid(True, color="#f0f0f0", linewidth=0.7)
    ax.xaxis.grid(False)
    ax.set_ylabel("Ratio %", fontsize=6, color="#6b7280", labelpad=1)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#9ca3af")
    ax.spines["bottom"].set_color("#9ca3af")
    fig.tight_layout(pad=0.25)

    output = io.BytesIO()
    canvas.print_png(output)
    output.seek(0)
    return output


def _cdf_points(values: List[float]) -> Tuple[List[float], List[float]]:
    sorted_values = sorted(values)
    count = len(sorted_values)
    if count == 0:
        return [], []
    return sorted_values, [((idx + 1) / count) * 100 for idx in range(count)]


def _histogram_series(values: List[float], bins: int, min_v: float, max_v: float, area_weighted: bool) -> List[float]:
    series = [0.0 for _ in range(bins)]
    if not values:
        return series
    span = max(max_v - min_v, 1e-6)
    total_weight = 0.0
    for raw in values:
        if area_weighted and raw <= 0:
            continue
        pct = (raw - min_v) / span
        idx = min(bins - 1, max(0, math.floor(pct * bins)))
        weight = math.pi * (raw / 2) * (raw / 2) if area_weighted else 1.0
        series[idx] += weight
        total_weight += weight
    if area_weighted and total_weight > 0:
        return [(value / total_weight) * 100 for value in series]
    return series


def _format_axis_tick(value: float) -> str:
    if abs(value) < 1e-6:
        return "0"
    if abs(value) >= 100 or abs(value - round(value)) < 1e-6:
        return str(int(round(value)))
    return f"{value:.1f}"


def _draw_cdf_series(slide, values: List[float], x: int, y: int, w: int, h: int, x_min: float, x_max: float, color: RGBColor):
    if not values:
        return
    sorted_values = sorted(values)
    points: List[Tuple[int, int]] = []
    span = max(x_max - x_min, 1e-6)
    count = len(sorted_values)
    for idx, value in enumerate(sorted_values):
        x_pct = min(1.0, max(0.0, (value - x_min) / span))
        y_pct = (idx + 1) / count
        points.append((int(x + x_pct * w), int(y + (1 - y_pct) * h)))
    _draw_polyline(slide, points, color, width=1.2)


def _draw_axes(slide, x: int, y: int, w: int, h: int, *, y_ticks: List[int]):
    axis_color = RGBColor(156, 163, 175)
    grid_color = RGBColor(229, 231, 235)
    for tick in y_ticks:
        ty = int(y + (1 - tick / max(y_ticks)) * h)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, ty, x + w, ty)
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.4)
    for x1, y1, x2, y2 in [(x, y, x, y + h), (x, y + h, x + w, y + h)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = axis_color
        line.line.width = Pt(0.8)


def _draw_polyline(slide, points: List[Tuple[int, int]], color: RGBColor, *, width: float = 1.0, dashed: bool = False):
    for start, end in zip(points, points[1:]):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start[0], start[1], end[0], end[1])
        line.line.color.rgb = color
        line.line.width = Pt(width)
        if dashed:
            line.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def _add_triangle_slide_legend(slide, x: float, y: float, selected_version_label: str, golden_version_label: str):
    _add_text(slide, "Golden 實線", x, y, 0.78, 0.16, size=7, color=RGBColor(75, 85, 99))
    solid = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.82), Inches(y + 0.08), Inches(x + 1.12), Inches(y + 0.08))
    solid.line.color.rgb = RGBColor(75, 85, 99)
    solid.line.width = Pt(1.2)
    _add_text(slide, golden_version_label, x + 1.18, y, 1.25, 0.16, size=7, color=RGBColor(75, 85, 99))

    _add_text(slide, "分析虛線", x + 2.55, y, 0.78, 0.16, size=7, color=RGBColor(75, 85, 99))
    dashed = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 3.37), Inches(y + 0.08), Inches(x + 3.67), Inches(y + 0.08))
    dashed.line.color.rgb = RGBColor(75, 85, 99)
    dashed.line.width = Pt(1.2)
    dashed.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _add_text(slide, selected_version_label, x + 3.73, y, 1.25, 0.16, size=7, color=RGBColor(75, 85, 99))

    cursor = x + 5.35
    for info in ROW_SERIES_INFO:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cursor), Inches(y + 0.035), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = info["color"]
        dot.line.fill.background()
        _add_text(slide, str(info["label"]), cursor + 0.12, y, 0.62, 0.16, size=7, color=RGBColor(75, 85, 99))
        cursor += 0.82


def _draw_triangle_panel(slide, x: float, y: float, w: float, h: float, col_key: str, dev: str, selected_snapshot: Snapshot, golden_snapshot: Snapshot):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(220, 223, 230)
    card.line.width = Pt(0.8)
    _add_text(slide, COL_LABELS[col_key], x + 0.15, y + 0.15, w - 0.3, 0.22, size=10, bold=True, color=RGBColor(0, 48, 180), align=PP_ALIGN.CENTER)

    cx = Inches(x + w / 2)
    cy = Inches(y + h * 0.47)
    radius = Inches(min(w, h) * 0.31)
    angles = [-math.pi / 2, math.pi / 6, 5 * math.pi / 6]
    max_pct = 50.0

    def point(pct: float, axis_idx: int) -> Tuple[int, int]:
        r = radius * (pct / max_pct)
        return (int(cx + r * math.cos(angles[axis_idx])), int(cy + r * math.sin(angles[axis_idx])))

    for level in [10, 25, 50]:
        _draw_polygon(slide, [point(level, i) for i in range(3)], RGBColor(209, 213, 219), width=0.55)
    for i in range(3):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, *point(50, i))
        line.line.color.rgb = RGBColor(156, 163, 175)
        line.line.width = Pt(0.55)
        lx, ly = point(58, i)
        _add_text(slide, f"[{ORIENT_LABELS[i]}]", lx / 914400 - 0.2, ly / 914400 - 0.08, 0.4, 0.14, size=7, bold=True, align=PP_ALIGN.CENTER)

    for info in ROW_SERIES_INFO:
        row_key = str(info["rowKey"])
        pos = f"{col_key}-{row_key}"
        color = info["color"]
        selected_values = _orient_values(selected_snapshot.get(pos), dev)
        golden_values = _orient_values(golden_snapshot.get(pos), dev)
        legend_y = y + h - 1.0 + ROW_KEYS.index(row_key) * 0.22
        _add_triangle_panel_legend_line(slide, x + 0.52, legend_y, f"{pos} Golden", color, dashed=False)
        _add_triangle_panel_legend_line(slide, x + 2.0, legend_y, f"{pos} Sample", color, dashed=True)
        if selected_values:
            sample_points = [point(v, i) for i, v in enumerate(selected_values)]
            _draw_polygon(slide, sample_points, color, width=1.0, dashed=True)
            _add_triangle_value_labels(slide, sample_points, selected_values, color, dashed=True)
        if golden_values:
            golden_points = [point(v, i) for i, v in enumerate(golden_values)]
            _draw_polygon(slide, golden_points, color, width=1.1, dashed=False)
            _add_triangle_value_labels(slide, golden_points, golden_values, color, dashed=False)


def _add_triangle_panel_legend_line(slide, x: float, y: float, text: str, color: RGBColor, *, dashed: bool):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y + 0.08), Inches(x + 0.24), Inches(y + 0.08))
    line.line.color.rgb = color
    line.line.width = Pt(1.0)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _add_text(slide, text, x + 0.28, y, 1.05, 0.16, size=6, color=RGBColor(75, 85, 99))


def _add_triangle_value_labels(slide, points: List[Tuple[int, int]], values: List[float], color: RGBColor, *, dashed: bool):
    for idx, ((px, py), value) in enumerate(zip(points, values)):
        dx = [-0.13, 0.04, -0.22][idx]
        dy = [-0.13, -0.02, -0.02][idx]
        if dashed:
            dy += 0.08
        _add_text(slide, f"{value:.1f}%", px / 914400 + dx, py / 914400 + dy, 0.3, 0.1, size=4, color=color, align=PP_ALIGN.CENTER)


def _draw_polygon(slide, points: List[Tuple[int, int]], color: RGBColor, *, width: float, dashed: bool = False):
    if len(points) < 2:
        return
    _draw_polyline(slide, points + [points[0]], color, width=width, dashed=dashed)


def _add_orientation_table(slide, x: float, y: float, w: float, h: float, col_key: str, dev: str, selected_snapshot: Snapshot, golden_snapshot: Snapshot):
    table = slide.shapes.add_table(4, 4, Inches(x), Inches(y), Inches(w), Inches(h)).table
    _style_header_cell(table.cell(0, 0), "位置")
    for c, label in enumerate(ORIENT_LABELS, start=1):
        _style_header_cell(table.cell(0, c), label)
    for r, row_key in enumerate(ROW_KEYS, start=1):
        pos = f"{col_key}-{row_key}"
        _style_header_cell(table.cell(r, 0), pos)
        selected = _orient_values(selected_snapshot.get(pos), dev) or [0.0, 0.0, 0.0]
        golden = _orient_values(golden_snapshot.get(pos), dev) or [0.0, 0.0, 0.0]
        for c in range(3):
            _set_cell_text(table.cell(r, c + 1), f"S {selected[c]:.1f}%\nG {golden[c]:.1f}%", size=7)


def _draw_orientation_line_chart(slide, x: float, y: float, w: float, h: float, title: str, row_key: str, dev: str, selected_snapshot: Snapshot):
    _add_text(slide, title, x, y + 0.47, 0.95, 0.2, size=10, bold=True)
    px = Inches(x + 1.05)
    py = Inches(y + 0.08)
    pw = Inches(w - 1.35)
    ph = Inches(h - 0.35)
    _draw_axes(slide, px, py, pw, ph, y_ticks=[0, 25, 50, 75, 100])
    for idx, tick in enumerate(ORIENT_LINE_X_TICKS):
        lx = x + 1.05 + (w - 1.35) * idx / 2
        _add_text(slide, str(tick["label"]), lx - 0.15, y + h - 0.18, 0.3, 0.12, size=7, color=RGBColor(107, 114, 128), align=PP_ALIGN.CENTER)
    for info in ORIENT_LINE_SERIES_INFO:
        orientation_index = int(info["index"])
        values = []
        for tick in ORIENT_LINE_X_TICKS:
            orient_values = _orient_values(selected_snapshot.get(f"{tick['colKey']}-{row_key}"), dev) or [0.0, 0.0, 0.0]
            values.append(orient_values[orientation_index])
        points = []
        for idx, value in enumerate(values):
            xx = int(px + pw * idx / 2)
            yy = int(py + ph * (1 - min(100.0, max(0.0, value)) / 100.0))
            points.append((xx, yy))
        _draw_polyline(slide, points, info["color"], width=1.5)
        for xx, yy in points:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, xx - Inches(0.035), yy - Inches(0.035), Inches(0.07), Inches(0.07))
            dot.fill.solid()
            dot.fill.fore_color.rgb = info["color"]
            dot.line.fill.background()


def _add_native_orientation_line_chart(slide, x: float, y: float, w: float, h: float, title: str, row_key: str, dev: str, selected_snapshot: Snapshot):
    _add_text(slide, title, x, y, 1.25, 0.22, size=10, bold=True)
    data = CategoryChartData()
    data.categories = [str(tick["label"]) for tick in ORIENT_LINE_X_TICKS]
    for info in ORIENT_LINE_SERIES_INFO:
        orientation_index = int(info["index"])
        values = []
        for tick in ORIENT_LINE_X_TICKS:
            orient_values = _orient_values(selected_snapshot.get(f"{tick['colKey']}-{row_key}"), dev) or [0.0, 0.0, 0.0]
            values.append(orient_values[orientation_index])
        data.add_series(str(info["label"]), values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(x + 1.25),
        Inches(y - 0.02),
        Inches(w - 1.25),
        Inches(h),
        data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.major_unit = 25
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    for idx, series in enumerate(chart.series):
        color = ORIENT_LINE_SERIES_INFO[idx]["color"]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(1.6)
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = color


def _add_legend(slide, x: float, y: float, items: List[Tuple[str, RGBColor]]):
    cursor = x
    for label, color in items:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cursor), Inches(y + 0.12), Inches(cursor + 0.28), Inches(y + 0.12))
        line.line.color.rgb = color
        line.line.width = Pt(2)
        _add_text(slide, label, cursor + 0.34, y, 2.25, 0.22, size=9, color=RGBColor(75, 85, 99))
        cursor += 2.55


def _add_placeholder(slide, x: float, y: float, w: float, h: float, text: str):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(243, 244, 246)
    rect.line.color.rgb = RGBColor(209, 213, 219)
    _add_text(slide, text, x, y + h / 2 - 0.1, w, 0.2, size=10, color=RGBColor(107, 114, 128), align=PP_ALIGN.CENTER)


def _format_nine_grid_cell(pos: str, features: Optional[FeatureMap]) -> str:
    if not features:
        return f"{pos}\nSample 未上傳"
    grains = _finite_list(features.get("grains", []))
    if grains:
        sorted_grains = sorted(grains)
        mean = statistics.mean(grains)
        max_v = sorted_grains[-1]
        p75 = sorted_grains[min(len(sorted_grains) - 1, int(len(sorted_grains) * 0.75))]
        grain_text = f"Grain Size (um)\nMean {mean:.2f} / Max {max_v:.2f}\nP75 {p75:.2f} / Count {len(grains)}"
    else:
        grain_text = "Grain Size (um)\n-"
    orient20 = _orient_values(features, "20%") or [0.0, 0.0, 0.0]
    orient15 = _orient_values(features, "15%") or [0.0, 0.0, 0.0]
    return (
        f"{pos}\n{grain_text}\n"
        f"Orientation Ratio\n20deg [001]{orient20[0]:.1f}% [011]{orient20[1]:.1f}% [111]{orient20[2]:.1f}%\n"
        f"15deg [001]{orient15[0]:.1f}% [011]{orient15[1]:.1f}% [111]{orient15[2]:.1f}%"
    )


def _ordered_positions() -> List[str]:
    return [f"{col_key}-{row_key}" for row_key in ROW_KEYS for col_key in COL_KEYS]


def _position_status(pos: str, features: Optional[FeatureMap], missing_golden: set, missing_selected: set) -> str:
    if not features or pos in missing_selected:
        return "Sample 未上傳"
    if pos in missing_golden:
        return "缺 Golden"
    return "OK"


def _grain_stats(features: Optional[FeatureMap]) -> Tuple[str, str, str, str]:
    grains = _finite_list(features.get("grains", [])) if features else []
    if not grains:
        return "-", "-", "-", "0"
    sorted_grains = sorted(grains)
    p75_idx = min(len(sorted_grains) - 1, math.floor(len(sorted_grains) * 0.75))
    return (
        f"{statistics.mean(grains):.2f}",
        f"{sorted_grains[-1]:.2f}",
        f"{sorted_grains[p75_idx]:.2f}",
        str(len(grains)),
    )


def _orient_values(features: Optional[FeatureMap], dev: str) -> Optional[List[float]]:
    if not features:
        return None
    values = features.get(f"orientation_ratio({dev})")
    if not isinstance(values, list):
        return None
    return [float(values[idx] if idx < len(values) else 0) * 100 for idx in range(3)]


def _finite_list(values: Any) -> List[float]:
    if not isinstance(values, list):
        return []
    return [float(v) for v in values if isinstance(v, (int, float)) and math.isfinite(float(v))]


def _sort_position_key(pos: str):
    col_order = {key: idx for idx, key in enumerate(COL_KEYS)}
    row_order = {key: idx for idx, key in enumerate(ROW_KEYS)}
    if "-" in pos:
        col, row = pos.split("-", 1)
        return (row_order.get(row, 99), col_order.get(col, 99), pos)
    return (99, 99, pos)
